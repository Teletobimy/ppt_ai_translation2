# translate_ppt.py  (서식 보존 + 톤 선택 버전)
# pyinstaller --onefile --name BNK_TranslatePPT PPT_Language_Change.py
import os
import re
import time
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openai


# ---------- [서식 보존을 위한 태깅/복원 유틸] ----------
RUN_TAG = re.compile(r"\[\[R(\d+)\]\]|\[\[/R(\d+)\]\]")
P_TAG = re.compile(r"\[\[P(\d+)\]\]|\[\[/P(\d+)\]\]")
# ==== [설정] =====================================================
# 맨 위 import 근처
SORRY_PATTERNS = [
    "i'm sorry", "i am sorry", "sorry, but", "apologize",
    "죄송하지만", "죄송합니다", "번역할 내용이 없습니다"
]

def is_effectively_empty_tagged(tagged_text: str) -> bool:
    """[[R#]]마커를 제거하고 남는 콘텐츠가 실질적으로 비었는지 판단"""
    stripped = RUN_TAG.sub("", tagged_text)  # 마커 제거
    return stripped.strip() == ""  # 공백만 남으면 빈 것으로 간주

def looks_like_apology(text: str) -> bool:
    low = (text or "").lower()
    return any(p in low for p in SORRY_PATTERNS)

# NOTE: API keys are no longer embedded. Provide them via environment variables:
# - OPENAI_API_KEY
# - DEEPSEEK_API_KEY (optional, for DeepSeek usage)

LANG_OPTIONS = [
    "English",
    "Indonesian",
    "Italian",
    "French",
    "Spanish",
    "Korean",
    "Japanese",
    "Russian",
    "German",
    "Portuguese",
    "Chinese (Simplified)",
    "Chinese (Traditional)",
]

# ✅ 톤 옵션 추가
TONE_OPTIONS = [
    "기본값",
    "Med/Pharma Pro (20y)",   # 의료기기/전문약사 20년 전문가
    "Beauty Pro (20y, chic)", # 세련된 뷰티 20년 전문가
    "GenZ Female (20s)",      # 20대 여성 타깃
]

OPENAI_MODEL = "gpt-4o"
DEEPSEEK_MODEL = "deepseek-chat"
SLEEP_SEC = 0
# ===============================================================


def unique_path(path: str) -> str:  
    if not os.path.exists(path):
        return path
    base, ext = os.path.splitext(path)
    i = 1
    while True:
        candidate = f"{base} ({i}){ext}"
        if not os.path.exists(candidate):
            return candidate
        i += 1

def create_deepseek_client():
    """DeepSeek 클라이언트 생성"""
    api_key = os.getenv("DEEPSEEK_API_KEY", "")
    if not api_key:
        raise RuntimeError("DEEPSEEK_API_KEY is not set. Configure it in secrets or environment.")
    return openai.OpenAI(api_key=api_key, base_url="https://api.deepseek.com")


def safe_request(client, prompt, retries=3, delay=3, use_deepseek=False):
    for attempt in range(retries):
        try:
            model = DEEPSEEK_MODEL if use_deepseek else OPENAI_MODEL
            resp = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                timeout=60,
            )
            content = ""
            if resp and hasattr(resp, "choices") and resp.choices:
                content = getattr(resp.choices[0].message, "content", "") or ""
            if content:
                return content.strip()
        except Exception as e:
            print(f"⚠️ Error (attempt {attempt+1}): {e}")
            with open("error.log", "a", encoding="utf-8") as f:
                f.write(f"[Attempt {attempt+1}] {e}\n")
            time.sleep(delay)
    return ""




def _font_attrs(run):
    f = run.font
    # 값 그대로 보존(None 포함)
    name = f.name                 # None이면 테마/마스터 상속
    size = f.size.pt if f.size else None
    bold = f.bold                 # True/False/None
    italic = f.italic             # True/False/None
    underline = f.underline       # True/False/None

    rgb = None
    if f.color is not None and getattr(f.color, "rgb", None) is not None:
        rgb = (f.color.rgb[0], f.color.rgb[1], f.color.rgb[2])

    return {"name": name, "size": size, "bold": bold, "italic": italic,
            "underline": underline, "rgb": rgb}

def _apply_font_attrs(run, attrs):
    from pptx.util import Pt
    f = run.font

    if attrs.get("name") is not None:
        f.name = attrs["name"]
    if attrs.get("size") is not None:
        f.size = Pt(attrs["size"])
    if attrs.get("bold") is not None:
        f.bold = attrs["bold"]
    if attrs.get("italic") is not None:
        f.italic = attrs["italic"]
    if attrs.get("underline") is not None:
        f.underline = attrs["underline"]
    if attrs.get("rgb") is not None:
        r, g, b = attrs["rgb"]
        f.color.rgb = RGBColor(r, g, b)


def tag_paragraph(paragraph):
    text_parts, style_map, idx = [], {}, 1
    for run in paragraph.runs:
        t = run.text or ""
        if not t:
            continue
        style_map[idx] = _font_attrs(run)
        text_parts.append(f"[[R{idx}]]{t}[[/R{idx}]]")
        idx += 1
    return "".join(text_parts), style_map

def rebuild_paragraph_from_tagged(paragraph, translated, style_map):
    while paragraph.runs:
        paragraph.runs[0]._r.getparent().remove(paragraph.runs[0]._r)

    tokens, pos = [], 0
    for m in RUN_TAG.finditer(translated):
        s, e = m.span()
        if s > pos:
            tokens.append(("text", translated[pos:s]))
        if m.group(1):
            tokens.append(("start", int(m.group(1))))
        if m.group(2):
            tokens.append(("end", int(m.group(2))))
        pos = e
    if pos < len(translated):
        tokens.append(("text", translated[pos:]))

    out_runs, stack, buffer = [], [], {}
    for ttype, value in tokens:
        if ttype == "start":
            stack.append(value)
            buffer.setdefault(value, [])
        elif ttype == "end":
            if stack and stack[-1] == value:
                stack.pop()
            joined = "".join(buffer.get(value, []))
            out_runs.append((value, joined))
            buffer[value] = []
        else:
            if stack:
                buffer[stack[-1]].append(value)
            else:
                out_runs.append((None, value))

    for run_id, buf in buffer.items():
        if buf:
            out_runs.append((run_id, "".join(buf)))

    if not out_runs:
        r = paragraph.add_run()
        r.text = translated
        return

    for run_id, txt in out_runs:
        if not txt:
            continue
        r = paragraph.add_run()
        r.text = txt
        if run_id and run_id in style_map:
            _apply_font_attrs(r, style_map[run_id])

def _parse_run_chunks(translated):
    # R# → 텍스트 매핑과, 마커 밖 텍스트 존재 여부를 판정
    ids = []
    chunks = {}
    stack = []
    buf = {}
    outside = []

    pos = 0
    for m in RUN_TAG.finditer(translated):
        s, e = m.span()
        if s > pos:
            if stack:
                buf.setdefault(stack[-1], []).append(translated[pos:s])
            else:
                outside.append(translated[pos:s])
        if m.group(1):  # [[R#]]
            rid = int(m.group(1)); stack.append(rid); ids.append(rid)
        if m.group(2):  # [[/R#]]
            rid = int(m.group(2))
            if stack and stack[-1] == rid:
                stack.pop()
                joined = "".join(buf.get(rid, []))
                chunks[rid] = joined
                buf[rid] = []
        pos = e
    if pos < len(translated):
        if stack:
            buf.setdefault(stack[-1], []).append(translated[pos:])
        else:
            outside.append(translated[pos:])

    # 닫히지 않은 버퍼 처리
    for rid, lst in buf.items():
        if lst:
            chunks[rid] = chunks.get(rid, "") + "".join(lst)

    has_outside = any(t.strip() for t in outside)
    return ids, chunks, has_outside

def try_inplace_update_paragraph(paragraph, translated):
    """마커가 1..N으로 정확히 존재하고, 마커 밖 텍스트가 없으면
    기존 runs에 텍스트만 주입하여 서식을 100% 유지한다."""
    ids, chunks, has_outside = _parse_run_chunks(translated)
    runs = [r for r in paragraph.runs if (r.text or "") != ""]
    N = len(runs)

    # 조건: 마커 밖 텍스트가 없어야 하고, R1..RN이 정확히 한 번씩 존재
    if has_outside or N == 0 or set(ids) != set(range(1, N+1)) or any(ids.count(i) != 1 for i in range(1, N+1)):
        return False

    for i, run in enumerate(runs, start=1):
        run.text = chunks.get(i, "")
    return True


# ---------- [프롬프트 빌더] ----------
def build_tone_instructions(tone: str) -> str:
    """
    선택한 톤에 맞는 지시문을 반환
    """
    if tone == "기본값":
        return (
            "Use a natural, professional beauty-industry tone localized to the target market. "
            "Keep terminology consistent with beauty marketing and professional skincare. "
            "Be clear and persuasive without hype; avoid overpromising."
        )
    if tone == "Med/Pharma Pro (20y)":
        return (
            "Use a formal, clinically precise B2B tone suitable for medical devices and professional pharmacists. "
            "Prioritize clarity, compliance, and evidence-based wording. Avoid hype. "
            "Prefer terminology used in regulatory, clinical, and professional settings."
        )
    if tone == "Beauty Pro (20y, chic)":
        return (
            "Use a refined, polished professional tone common in premium beauty and aesthetic clinics. "
            "Balance expertise with approachable elegance. Maintain brand voice consistency without overpromising."
        )
    if tone == "GenZ Female (20s)":
        return (
            "Use a modern, friendly, and concise tone tailored for women in their 20s. "
            "Be clear and engaging for social content, but avoid slang overload, emojis, and exaggerated claims."
        )
    # fallback
    return "Use a neutral professional tone appropriate for the beauty industry."

def build_chinese_prompt(tagged_text: str, target_lang: str) -> str:
    """
    전문적인 한국어→중국어 번역을 위한 프롬프트 (간체/번체 구분)
    """
    chinese_type = "간체" if "Simplified" in target_lang else "번체"
    
    return (
        f"당신은 한국어를 정확하고 자연스러운 중국어({chinese_type})로 번역하는 전문가입니다.\n\n"
        f"# 주요 특징\n"
        f"- 프레젠테이션, 보고서, 비즈니스 문서 등에 적합한 공식적이고 세련된 표현 사용\n"
        f"- 문맥을 고려하여 문장의 의미와 뉘앙스를 세밀하게 분석하여 적절한 표현으로 번역\n"
        f"- 원문의 의도와 어조를 유지하되, 중국 원어민이 자연스럽게 들리도록 문장 구조 조정\n"
        f"- 번역 이외의 불필요한 설명 금지\n"
        f"- 창의적 재해석 없이 원문에 충실한 번역 수행\n"
        f"- 반드시 중국어 {chinese_type}로 번역하세요\n\n"
        f"# 고유명사 처리 규칙\n"
        f"- '피더린'은 'PYDERIN'으로 번역하세요 (브랜드명이므로 대문자로)\n"
        f"- 기타 고유명사(인명, 지명, 회사명, 브랜드명 등)는 번역하지 말고 원문 그대로 유지하세요\n"
        f"- 영어 고유명사는 그대로 유지하세요\n\n"
        f"다음 한국어 텍스트를 자연스러운 중국어({chinese_type})로 번역하세요.\n"
        f"중요: 단락 마커 [[P#]]...[[/P#]]와 런 마커 [[R#]]...[[/R#]]는 절대 변경하거나 제거하지 마세요.\n"
        f"- 단락 개수(P#)와 순서를 정확히 유지하세요.\n"
        f"- 마커 내부 텍스트만 번역하고, 마커 자체는 그대로 두세요.\n\n"
        f"번역할 텍스트:\n{tagged_text}"
    )

def build_prompt(tagged_text: str, target_lang: str, tone: str) -> str:
    # Chinese translation uses specialized prompt
    if "Chinese" in target_lang:
        return build_chinese_prompt(tagged_text, target_lang)
    
    tone_text = build_tone_instructions(tone)
    return (
        f"Translate the following beauty industry presentation text into natural, professional {target_lang}. "
        f"Only return the translated text. If there is nothing to translate, return an empty string. "
        f"Context: {tone_text} "
        f"Avoid literal translation—use expressions that sound natural for beauty marketing and professional skincare. "
        f"If the source is already in {target_lang}, lightly copyedit for clarity, consistency, and terminology. "
        f"CRITICAL: Do NOT alter or remove any marker tags. Preserve both paragraph markers [[P#]]...[[/P#]] and run markers [[R#]]...[[/R#]] exactly, including counts and order. "
        f"Return ONLY the translated text with all markers preserved:\n\n{tagged_text}"
    )



# ---------- [번역 호출] ----------
def gpt_translate_tagged(tagged_text: str, client, target_lang: str, tone: str, use_deepseek=False) -> str:
    # 진짜 내용이 없으면 번역 스킵
    if not tagged_text.strip() or is_effectively_empty_tagged(tagged_text):
        return ""

    # 중국어 번역의 경우 DeepSeek 사용
    if "Chinese" in target_lang and use_deepseek:
        deepseek_client = create_deepseek_client()
        prompt = build_chinese_prompt(tagged_text, target_lang)
        content = safe_request(deepseek_client, prompt, retries=3, delay=3, use_deepseek=True)
    else:
        prompt = build_prompt(tagged_text, target_lang, tone)
        content = safe_request(client, prompt, retries=3, delay=3)

    # 실패 시 원문(마커 포함) 반환 → 원문 유지
    if not content:
        return tagged_text

    # 사과문/에러문구가 들어오면 원문 유지
    if looks_like_apology(content):
        return tagged_text

    return content

def gpt_review_chinese_translation(original_korean: str, translated_chinese: str, client, use_deepseek=False) -> dict:
    """
    중국어 번역의 자연스러움을 검토하고 필요시 수정된 번역을 반환
    """
    if not original_korean.strip() or not translated_chinese.strip():
        return {"is_awkward": False, "revised_translation": translated_chinese}
    
    review_prompt = (
        f"당신은 중국어 번역 품질을 검토하는 전문가입니다.\n\n"
        f"원문 (한국어): {original_korean}\n"
        f"번역문 (중국어): {translated_chinese}\n\n"
        f"다음을 검토해주세요:\n"
        f"1. 중국 원어민이 읽었을 때 어색하거나 부자연스러운 부분이 있는가?\n"
        f"2. 문법적으로 올바른가?\n"
        f"3. 표현이 자연스러운가?\n\n"
        f"응답 형식:\n"
        f"어색함: [YES/NO]\n"
        f"수정된 번역: [수정된 중국어 번역 또는 원래 번역]\n"
        f"설명: [어색한 이유 또는 수정 사항]\n\n"
        f"중요: [[P#]]와 [[R#]] 마커는 절대 변경하지 마세요."
    )
    
    try:
        if use_deepseek:
            deepseek_client = create_deepseek_client()
            content = safe_request(deepseek_client, review_prompt, retries=2, delay=2, use_deepseek=True)
        else:
            content = safe_request(client, review_prompt, retries=2, delay=2)
            
        if not content:
            return {"is_awkward": False, "revised_translation": translated_chinese}
        
        # Parse response
        lines = content.strip().split('\n')
        is_awkward = False
        revised_translation = translated_chinese
        
        for line in lines:
            if line.startswith("어색함:"):
                is_awkward = "YES" in line.upper()
            elif line.startswith("수정된 번역:"):
                revised_translation = line.replace("수정된 번역:", "").strip()
        
        return {"is_awkward": is_awkward, "revised_translation": revised_translation}
        
    except Exception as e:
        print(f"⚠️ Review error: {e}")
        return {"is_awkward": False, "revised_translation": translated_chinese}




# ---------- [파일/언어/톤 선택 UI] ----------
def choose_pptx_with_dialog() -> str:
    # Lazy import Tkinter for desktop-only usage
    import tkinter as tk  # type: ignore
    from tkinter import filedialog  # type: ignore

    root = tk.Tk()
    root.withdraw()
    root.update_idletasks()
    filepath = filedialog.askopenfilename(
        title="번역할 PPTX 파일 선택",
        filetypes=[("PowerPoint files", "*.pptx")],
    )
    root.destroy()
    return filepath or ""

def choose_language_with_window() -> str:
    sel = {"value": ""}

    def on_start():
        v = var.get().strip()
        if not v:
            from tkinter import messagebox  # type: ignore
            messagebox.showwarning("알림", "언어를 선택하세요.")
            return
        sel["value"] = v
        win.destroy()

    import tkinter as tk  # type: ignore
    win = tk.Tk()
    win.title("Target Language")
    win.geometry("360x160")
    win.resizable(False, False)

    frm = tk.Frame(win, padx=12, pady=12)
    frm.pack(fill="both", expand=True)

    tk.Label(frm, text="번역 대상 언어 선택:").pack(anchor="w", pady=(0, 6))

    var = tk.StringVar(value=LANG_OPTIONS[0])
    opt = tk.OptionMenu(frm, var, *LANG_OPTIONS)
    opt.pack(fill="x")

    tk.Button(frm, text="다음(톤 선택)", command=on_start).pack(pady=12)

    win.lift(); win.attributes("-topmost", True); win.after(200, lambda: win.attributes("-topmost", False))
    win.mainloop()
    return sel["value"]

def choose_tone_with_window(selected_language: str) -> tuple:
    sel = {"value": "", "use_deepseek": False}

    def on_start():
        v = var.get().strip()
        if not v:
            from tkinter import messagebox  # type: ignore
            messagebox.showwarning("알림", "톤을 선택하세요.")
            return
        sel["value"] = v
        sel["use_deepseek"] = deepseek_var.get()
        win.destroy()

    import tkinter as tk  # type: ignore
    win = tk.Tk()
    win.title("Target Tone & DeepSeek Option")
    win.geometry("450x280")
    win.resizable(False, False)

    frm = tk.Frame(win, padx=12, pady=12)
    frm.pack(fill="both", expand=True)

    tk.Label(frm, text="번역 톤 선택:").pack(anchor="w", pady=(0, 6))

    var = tk.StringVar(value=TONE_OPTIONS[0])
    opt = tk.OptionMenu(frm, var, *TONE_OPTIONS)
    opt.pack(fill="x")

    # DeepSeek 사용 옵션 (중국어일 때만 표시)
    if "Chinese" in selected_language:
        deepseek_frame = tk.Frame(frm)
        deepseek_frame.pack(fill="x", pady=(10, 0))
        
        deepseek_var = tk.BooleanVar(value=True)
        chinese_type = "간체" if "Simplified" in selected_language else "번체"
        deepseek_check = tk.Checkbutton(
            deepseek_frame,
            text=f"중국어({chinese_type}) 번역 시 DeepSeek 사용 (권장)",
            variable=deepseek_var,
            font=("Arial", 9, "bold")
        )
        deepseek_check.pack(anchor="w")
        
        deepseek_info = tk.Label(
            deepseek_frame,
            text=f"✓ DeepSeek은 중국어({chinese_type}) 번역에 특화되어 더 자연스러운 번역 결과를 제공합니다",
            font=("Arial", 8),
            fg="blue"
        )
        deepseek_info.pack(anchor="w", pady=(2, 0))
    else:
        deepseek_var = tk.BooleanVar(value=False)

    # 간단한 설명 라벨
    info = tk.Label(
        frm,
        justify="left",
        text=(
            "- 기본값: 일반뷰티업계, 직역 최대한 회피\n"
            "- Med/Pharma Pro: 의료기기/전문약사 20년 전문가 톤\n"
            "- Beauty Pro (chic): 프리미엄 뷰티 전문가 톤\n"
            "- GenZ Female: 20대 여성 타깃의 친근한 톤(과장·슬랭 과다 금지)"
        ),
    )
    info.pack(anchor="w", pady=8)

    tk.Button(frm, text="번역 시작", command=on_start).pack(pady=6)

    win.lift(); win.attributes("-topmost", True); win.after(200, lambda: win.attributes("-topmost", False))
    win.mainloop()
    return sel["value"], sel["use_deepseek"]


def choose_font_scale_window() -> int:
    """폰트 스케일(%) 선택 창. 기본 100."""
    sel = {"value": 100}

    def on_ok():
        try:
            v = int(entry.get().strip())
            if v < 50 or v > 300:
                from tkinter import messagebox  # type: ignore
                messagebox.showwarning("알림", "50% ~ 300% 사이의 값을 입력하세요.")
                return
            sel["value"] = v
            win.destroy()
        except Exception:
            from tkinter import messagebox  # type: ignore
            messagebox.showwarning("알림", "정수 % 값을 입력하세요 (예: 90, 100, 120)")

    import tkinter as tk  # type: ignore
    win = tk.Tk()
    win.title("Font Scale (%)")
    win.geometry("360x160")
    win.resizable(False, False)

    frm = tk.Frame(win, padx=12, pady=12)
    frm.pack(fill="both", expand=True)

    tk.Label(frm, text="번역 후 폰트 크기 배율(%)").pack(anchor="w", pady=(0, 6))

    quick = tk.Frame(frm)
    quick.pack(anchor="w", pady=(0, 6))

    def set_quick(val):
        entry.delete(0, tk.END)
        entry.insert(0, str(val))

    for val in (80, 90, 100, 110, 120, 130):
        tk.Button(quick, text=f"{val}%", command=lambda v=val: set_quick(v), width=6).pack(side="left", padx=2)

    entry = tk.Entry(frm)
    entry.insert(0, "100")
    entry.pack(fill="x")

    tk.Button(frm, text="확인", command=on_ok).pack(pady=10)

    win.lift(); win.attributes("-topmost", True); win.after(200, lambda: win.attributes("-topmost", False))
    win.mainloop()
    return sel["value"]


# ---------- [본 처리] ----------
def translate_presentation(pptx_path: str, target_lang: str, tone: str, use_deepseek=False, font_scale_percent: int = 100):
    print(f"📂 파일: {pptx_path}")
    print(f"🌐 대상 언어: {target_lang}")
    print(f"🎙 톤: {tone}")
    if target_lang == "Chinese" and use_deepseek:
        print("🤖 DeepSeek 모델 사용 중...")
    else:
        print("🔑 OpenAI 클라이언트 초기화 중...")

    # Create OpenAI client from environment (Streamlit will provide via st.secrets)
    openai_key = os.getenv("OPENAI_API_KEY", "")
    if not openai_key:
        raise RuntimeError("OPENAI_API_KEY is not set. Configure it in secrets or environment.")
    client = openai.OpenAI(api_key=openai_key)

    print("📖 프레젠테이션 로딩...")
    pres = Presentation(pptx_path)

    slide_count = len(pres.slides)
    print(f"🖼 슬라이드 수: {slide_count}")
    print(f"🔍 폰트 배율: {font_scale_percent}%")
    
    # ---------- [블록 태깅/재구성 유틸] ----------
    def tag_paragraphs_block(paragraphs):
        """단락 단위로 [[P#]] 래핑하고, 각 단락 내부는 기존 [[R#]] 마커로 태깅.
        반환: (block_text, per_para_style_maps)
        """
        parts = []
        style_maps = []
        para_index = 1
        has_any = False
        for p in paragraphs:
            tagged, style_map = tag_paragraph(p)
            style_maps.append(style_map)
            if tagged:
                has_any = True
            parts.append(f"[[P{para_index}]]{tagged}[[/P{para_index}]]")
            para_index += 1
        return ("".join(parts), style_maps, has_any)

    def _parse_p_blocks(translated_block: str):
        """블록 텍스트에서 P 마커에 해당하는 텍스트를 추출.
        반환: (ids_in_order, p_to_inner_text, has_outside)
        """
        ids = []
        p_chunks = {}
        stack = []
        buf = {}
        outside = []

        pos = 0
        for m in P_TAG.finditer(translated_block):
            s, e = m.span()
            if s > pos:
                if stack:
                    buf.setdefault(stack[-1], []).append(translated_block[pos:s])
                else:
                    outside.append(translated_block[pos:s])
            if m.group(1):  # [[P#]]
                pid = int(m.group(1)); stack.append(pid); ids.append(pid)
            if m.group(2):  # [[/P#]]
                pid = int(m.group(2))
                if stack and stack[-1] == pid:
                    stack.pop()
                    joined = "".join(buf.get(pid, []))
                    p_chunks[pid] = joined
                    buf[pid] = []
            pos = e
        if pos < len(translated_block):
            if stack:
                buf.setdefault(stack[-1], []).append(translated_block[pos:])
            else:
                outside.append(translated_block[pos:])

        for pid, lst in buf.items():
            if lst:
                p_chunks[pid] = p_chunks.get(pid, "") + "".join(lst)

        has_outside = any(t.strip() for t in outside)
        return ids, p_chunks, has_outside

    def rebuild_block_from_tagged(paragraphs, translated_block: str, style_maps):
        """P 블록을 파싱하여 각 단락에 대해 기존 런을 보존하며 텍스트 주입.
        마커 또는 개수 불일치 시 False 반환.
        """
        ids, p_chunks, has_outside = _parse_p_blocks(translated_block)
        N = len(paragraphs)
        if has_outside or N == 0:
            return False
        # P1..PN이 정확히 1회씩 존재하는지
        if set(ids) != set(range(1, N+1)) or any(ids.count(i) != 1 for i in range(1, N+1)):
            return False

        for i, p in enumerate(paragraphs, start=1):
            inner = p_chunks.get(i, "")
            # 먼저 인플레이스 시도 (런 보존)
            if not try_inplace_update_paragraph(p, inner):
                rebuild_paragraph_from_tagged(p, inner, style_maps[i-1])
        return True
    
    print("-" * 60)

    def process_paragraphs_block(paragraphs, chinese_review_enabled=False):
        block_tagged, style_maps, has_any = tag_paragraphs_block(paragraphs)
        if not has_any:
            return  # nothing to translate
        translated_block = gpt_translate_tagged(block_tagged, client, target_lang, tone, use_deepseek)
        translated_block = translated_block.strip().strip('"').strip("'")
        if chinese_review_enabled:
            review_result = gpt_review_chinese_translation(block_tagged, translated_block, client, use_deepseek)
            if review_result.get("is_awkward"):
                translated_block = review_result.get("revised_translation", translated_block)
                chinese_type = "간체" if "Simplified" in target_lang else "번체"
                print(f"   ✅ 어색한 번역 감지 (블록, 중국어 {chinese_type}) - 수정됨")
        # 우선 블록 단위 복원 시도
        if not rebuild_block_from_tagged(paragraphs, translated_block, style_maps):
            # P 마커 불일치 → 단락 단위로 폴백
            for p in paragraphs:
                tagged, style_map = tag_paragraph(p)
                if not tagged:
                    continue
                t = gpt_translate_tagged(tagged, client, target_lang, tone, use_deepseek)
                t = t.strip().strip('"').strip("'")
                if chinese_review_enabled:
                    rr = gpt_review_chinese_translation(tagged, t, client, use_deepseek)
                    if rr.get("is_awkward"):
                        t = rr.get("revised_translation", t)
                if not try_inplace_update_paragraph(p, t):
                    rebuild_paragraph_from_tagged(p, t, style_map)
                time.sleep(SLEEP_SEC)

    def traverse_shape(shape):
        # 그룹: 재귀
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            for inner in shape.shapes:
                traverse_shape(inner)
            return

        # 텍스트 프레임
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            tf = shape.text_frame
            process_paragraphs_block(tf.paragraphs, chinese_review_enabled=("Chinese" in target_lang))
            return

        # 표
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if getattr(cell, "text_frame", None):
                        process_paragraphs_block(cell.text_frame.paragraphs, chinese_review_enabled=("Chinese" in target_lang))
            return

    for s_idx, slide in enumerate(pres.slides, start=1):
        print(f"▶ 슬라이드 {s_idx}/{slide_count}")
        for shape in slide.shapes:
            traverse_shape(shape)

    # ---------- [폰트 스케일 적용] ----------
    def apply_font_scale(presentation, scale_percent: int):
        from pptx.util import Pt
        factor = max(1, scale_percent) / 100.0

        def scale_run(run):
            if run.font.size is not None:
                try:
                    run.font.size = Pt(run.font.size.pt * factor)
                except Exception:
                    pass

        def scale_paragraph(paragraph):
            # 문단 레벨 폰트가 지정된 경우 스케일
            if paragraph.font is not None and paragraph.font.size is not None:
                try:
                    paragraph.font.size = Pt(paragraph.font.size.pt * factor)
                except Exception:
                    pass
            # 각 런 스케일
            for r in paragraph.runs:
                scale_run(r)

        def traverse_scale_shape(shape):
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                for inner in shape.shapes:
                    traverse_scale_shape(inner)
                return
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    scale_paragraph(p)
                return
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if getattr(cell, "text_frame", None):
                            for p in cell.text_frame.paragraphs:
                                scale_paragraph(p)
                return

        if abs(factor - 1.0) < 1e-6:
            return
        for slide in presentation.slides:
            for shape in slide.shapes:
                traverse_scale_shape(shape)

    apply_font_scale(pres, font_scale_percent)

    folder = os.path.dirname(pptx_path)
    stem, ext = os.path.splitext(os.path.basename(pptx_path))
    
    # 중국어 번역의 경우 톤 대신 중국어 타입을 사용
    if "Chinese" in target_lang:
        chinese_type = "Simplified" if "Simplified" in target_lang else "Traditional"
        outfile_name = f"{stem}_Chinese_{chinese_type}_AI번역완료{ext}"
    else:
        safe_tone = re.sub(r'[^0-9A-Za-z가-힣_()+-]', '', tone.replace(' ', ''))
        outfile_name = f"{stem}_{target_lang}_{safe_tone}_AI번역완료{ext}"
    
    outfile_path = os.path.join(folder, outfile_name)
    outfile_path = unique_path(outfile_path)

    print("-" * 60)
    print("💾 저장 중...")
    pres.save(outfile_path)
    print(f"✅ 번역 완료! 저장된 파일: {outfile_path}")
    return outfile_path


def main():
    pptx_path = choose_pptx_with_dialog()
    if not pptx_path:
        print("❌ 파일을 선택하지 않았습니다. 종료합니다.")
        return

    target_lang = choose_language_with_window()
    if not target_lang:
        print("❌ 언어를 선택하지 않았습니다. 종료합니다.")
        return

    tone, use_deepseek = choose_tone_with_window(target_lang)
    if not tone:
        print("❌ 톤을 선택하지 않았습니다. 종료합니다.")
        return
    font_scale = choose_font_scale_window()
    translate_presentation(pptx_path, target_lang, tone, use_deepseek, font_scale_percent=font_scale)


if __name__ == "__main__":
    main()
